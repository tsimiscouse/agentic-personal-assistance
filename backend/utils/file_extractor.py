"""
File Extractor Utility
Extracts text from various file formats:
- PDF (.pdf)
- Word Documents (.docx)
- PowerPoint Presentations (.pptx)
- Excel Files (.xlsx, .xls)
- CSV Files (.csv)
"""

import io
import base64
from typing import Optional, Dict
from loguru import logger

# PDF Processing
import PyPDF2

# Document Processing
try:
    import docx  # python-docx
except ImportError:
    docx = None
    logger.warning("python-docx not installed. DOCX support disabled.")

try:
    from pptx import Presentation  # python-pptx
except ImportError:
    Presentation = None
    logger.warning("python-pptx not installed. PPTX support disabled.")

try:
    import openpyxl
    import pandas as pd
except ImportError:
    openpyxl = None
    pd = None
    logger.warning("openpyxl or pandas not installed. Excel/CSV support disabled.")


# ============================================
# FILE TYPE DETECTION
# ============================================

def detect_file_type(filename: str) -> str:
    """
    Detect file type from filename extension.

    Args:
        filename: Name of the file

    Returns:
        str: File type (pdf, docx, pptx, xlsx, xls, csv, txt)
    """
    filename_lower = filename.lower()

    if filename_lower.endswith('.pdf'):
        return 'pdf'
    elif filename_lower.endswith('.docx'):
        return 'docx'
    elif filename_lower.endswith('.doc'):
        return 'doc'  # Legacy Word format
    elif filename_lower.endswith('.pptx'):
        return 'pptx'
    elif filename_lower.endswith('.ppt'):
        return 'ppt'  # Legacy PowerPoint format
    elif filename_lower.endswith('.xlsx'):
        return 'xlsx'
    elif filename_lower.endswith('.xls'):
        return 'xls'
    elif filename_lower.endswith('.csv'):
        return 'csv'
    elif filename_lower.endswith('.txt'):
        return 'txt'
    else:
        return 'unknown'


# ============================================
# MAIN EXTRACTION FUNCTION
# ============================================

def extract_text_from_file(file_data: bytes, filename: str, max_chars: int = 15000) -> Dict:
    """
    Extract text from file data.

    Args:
        file_data: File content as bytes
        filename: Original filename (for type detection)
        max_chars: Maximum characters to extract (for efficiency)

    Returns:
        Dict: {
            'success': bool,
            'text': str,
            'file_type': str,
            'metadata': dict,
            'error': str (optional)
        }
    """
    try:
        file_type = detect_file_type(filename)
        logger.info(f"Extracting text from {file_type} file: {filename}")

        # Route to appropriate extractor
        if file_type == 'pdf':
            result = _extract_from_pdf(file_data, max_chars)
        elif file_type == 'docx':
            result = _extract_from_docx(file_data, max_chars)
        elif file_type == 'pptx':
            result = _extract_from_pptx(file_data, max_chars)
        elif file_type in ['xlsx', 'xls']:
            result = _extract_from_excel(file_data, file_type, max_chars)
        elif file_type == 'csv':
            result = _extract_from_csv(file_data, max_chars)
        elif file_type == 'txt':
            result = _extract_from_txt(file_data, max_chars)
        else:
            return {
                'success': False,
                'text': '',
                'file_type': file_type,
                'error': f'Unsupported file type: {file_type}'
            }

        # Add file type to result
        result['file_type'] = file_type
        return result

    except Exception as e:
        logger.error(f"Error extracting text from file: {e}")
        return {
            'success': False,
            'text': '',
            'file_type': 'unknown',
            'error': str(e)
        }


# ============================================
# PDF EXTRACTION
# ============================================

def _extract_from_pdf(file_data: bytes, max_chars: int) -> Dict:
    """Extract text from PDF file."""
    try:
        pdf_file = io.BytesIO(file_data)
        pdf_reader = PyPDF2.PdfReader(pdf_file)

        num_pages = len(pdf_reader.pages)
        logger.info(f"PDF has {num_pages} pages")

        extracted_text = []
        total_chars = 0

        # Extract from each page until max_chars reached
        for page_num in range(num_pages):
            if total_chars >= max_chars:
                break

            page = pdf_reader.pages[page_num]
            text = page.extract_text()

            if text:
                extracted_text.append(text)
                total_chars += len(text)

        full_text = '\n'.join(extracted_text)

        # Truncate if needed
        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Content truncated for processing efficiency]"

        logger.info(f"Extracted {len(full_text)} characters from PDF")

        return {
            'success': True,
            'text': full_text,
            'metadata': {
                'pages': num_pages,
                'chars_extracted': len(full_text)
            }
        }

    except Exception as e:
        logger.error(f"Error extracting from PDF: {e}")
        return {
            'success': False,
            'text': '',
            'error': str(e)
        }


# ============================================
# WORD DOCUMENT EXTRACTION
# ============================================

def _extract_from_docx(file_data: bytes, max_chars: int) -> Dict:
    """Extract text from Word document (.docx)."""
    try:
        if docx is None:
            return {
                'success': False,
                'text': '',
                'error': 'python-docx library not installed'
            }

        doc_file = io.BytesIO(file_data)
        doc = docx.Document(doc_file)

        extracted_text = []
        total_chars = 0

        # Extract from paragraphs
        for paragraph in doc.paragraphs:
            if total_chars >= max_chars:
                break

            text = paragraph.text.strip()
            if text:
                extracted_text.append(text)
                total_chars += len(text)

        # Extract from tables
        for table in doc.tables:
            if total_chars >= max_chars:
                break

            for row in table.rows:
                row_text = ' | '.join([cell.text.strip() for cell in row.cells])
                if row_text:
                    extracted_text.append(row_text)
                    total_chars += len(row_text)

        full_text = '\n'.join(extracted_text)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Content truncated]"

        logger.info(f"Extracted {len(full_text)} characters from DOCX")

        return {
            'success': True,
            'text': full_text,
            'metadata': {
                'paragraphs': len(doc.paragraphs),
                'tables': len(doc.tables),
                'chars_extracted': len(full_text)
            }
        }

    except Exception as e:
        logger.error(f"Error extracting from DOCX: {e}")
        return {
            'success': False,
            'text': '',
            'error': str(e)
        }


# ============================================
# POWERPOINT EXTRACTION
# ============================================

def _extract_from_pptx(file_data: bytes, max_chars: int) -> Dict:
    """Extract text from PowerPoint presentation (.pptx)."""
    try:
        if Presentation is None:
            return {
                'success': False,
                'text': '',
                'error': 'python-pptx library not installed'
            }

        ppt_file = io.BytesIO(file_data)
        prs = Presentation(ppt_file)

        extracted_text = []
        total_chars = 0
        slide_count = 0

        # Extract from each slide
        for slide in prs.slides:
            if total_chars >= max_chars:
                break

            slide_count += 1
            slide_text = [f"\n--- Slide {slide_count} ---"]

            # Extract from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                        total_chars += len(text)

            if len(slide_text) > 1:  # More than just the header
                extracted_text.extend(slide_text)

        full_text = '\n'.join(extracted_text)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Content truncated]"

        logger.info(f"Extracted {len(full_text)} characters from PPTX ({slide_count} slides)")

        return {
            'success': True,
            'text': full_text,
            'metadata': {
                'slides': len(prs.slides),
                'processed_slides': slide_count,
                'chars_extracted': len(full_text)
            }
        }

    except Exception as e:
        logger.error(f"Error extracting from PPTX: {e}")
        return {
            'success': False,
            'text': '',
            'error': str(e)
        }


# ============================================
# EXCEL EXTRACTION
# ============================================

def _extract_from_excel(file_data: bytes, file_type: str, max_chars: int) -> Dict:
    """Extract text from Excel file (.xlsx, .xls)."""
    try:
        if pd is None:
            return {
                'success': False,
                'text': '',
                'error': 'pandas/openpyxl library not installed'
            }

        excel_file = io.BytesIO(file_data)

        # Read Excel file
        if file_type == 'xlsx':
            excel_data = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
        else:  # xls
            excel_data = pd.read_excel(excel_file, sheet_name=None, engine='xlrd')

        extracted_text = []
        total_chars = 0

        # Process each sheet
        for sheet_name, df in excel_data.items():
            if total_chars >= max_chars:
                break

            extracted_text.append(f"\n--- Sheet: {sheet_name} ---")

            # Convert dataframe to string representation
            sheet_text = df.to_string(index=False, max_rows=100)
            extracted_text.append(sheet_text)
            total_chars += len(sheet_text)

            # Add summary
            summary = f"\nRows: {len(df)}, Columns: {len(df.columns)}"
            extracted_text.append(summary)

        full_text = '\n'.join(extracted_text)

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Content truncated]"

        logger.info(f"Extracted {len(full_text)} characters from Excel ({len(excel_data)} sheets)")

        return {
            'success': True,
            'text': full_text,
            'metadata': {
                'sheets': len(excel_data),
                'chars_extracted': len(full_text)
            }
        }

    except Exception as e:
        logger.error(f"Error extracting from Excel: {e}")
        return {
            'success': False,
            'text': '',
            'error': str(e)
        }


# ============================================
# CSV EXTRACTION
# ============================================

def _extract_from_csv(file_data: bytes, max_chars: int) -> Dict:
    """Extract text from CSV file."""
    try:
        if pd is None:
            return {
                'success': False,
                'text': '',
                'error': 'pandas library not installed'
            }

        csv_file = io.BytesIO(file_data)
        df = pd.read_csv(csv_file)

        # Convert to string with max 200 rows for efficiency
        csv_text = df.to_string(index=False, max_rows=200)

        # Add summary
        summary = f"\n\nData Summary:\nRows: {len(df)}\nColumns: {len(df.columns)}\nColumn Names: {', '.join(df.columns.tolist())}"

        full_text = csv_text + summary

        if len(full_text) > max_chars:
            full_text = full_text[:max_chars] + "\n\n[Content truncated]"

        logger.info(f"Extracted {len(full_text)} characters from CSV")

        return {
            'success': True,
            'text': full_text,
            'metadata': {
                'rows': len(df),
                'columns': len(df.columns),
                'chars_extracted': len(full_text)
            }
        }

    except Exception as e:
        logger.error(f"Error extracting from CSV: {e}")
        return {
            'success': False,
            'text': '',
            'error': str(e)
        }


# ============================================
# TEXT FILE EXTRACTION
# ============================================

def _extract_from_txt(file_data: bytes, max_chars: int) -> Dict:
    """Extract text from plain text file."""
    try:
        # Try different encodings
        encodings = ['utf-8', 'latin-1', 'cp1252']
        text = None

        for encoding in encodings:
            try:
                text = file_data.decode(encoding)
                break
            except UnicodeDecodeError:
                continue

        if text is None:
            return {
                'success': False,
                'text': '',
                'error': 'Could not decode text file'
            }

        if len(text) > max_chars:
            text = text[:max_chars] + "\n\n[Content truncated]"

        logger.info(f"Extracted {len(text)} characters from TXT")

        return {
            'success': True,
            'text': text,
            'metadata': {
                'chars_extracted': len(text)
            }
        }

    except Exception as e:
        logger.error(f"Error extracting from TXT: {e}")
        return {
            'success': False,
            'text': '',
            'error': str(e)
        }


# ============================================
# BASE64 HELPER FUNCTIONS
# ============================================

def extract_text_from_base64(base64_data: str, filename: str, max_chars: int = 15000) -> Dict:
    """
    Extract text from base64 encoded file data.

    Args:
        base64_data: Base64 encoded file content
        filename: Original filename
        max_chars: Maximum characters to extract

    Returns:
        Dict: Extraction result
    """
    try:
        # Decode base64
        file_data = base64.b64decode(base64_data)

        # Extract text
        return extract_text_from_file(file_data, filename, max_chars)

    except Exception as e:
        logger.error(f"Error decoding base64: {e}")
        return {
            'success': False,
            'text': '',
            'file_type': 'unknown',
            'error': f'Base64 decode error: {str(e)}'
        }
